import io
from typing import Iterable

from docx import Document
from pptx import Presentation

from app.models import DocumentStructure, Project


def export_docx(project: Project, structures: Iterable[DocumentStructure]) -> bytes:
    doc = Document()
    doc.add_heading(project.project_name, level=1)
    for structure in sorted(structures, key=lambda s: s.order_index):
        doc.add_heading(structure.title, level=2)
        body = structure.content.generated_content if structure.content else "Content pending."
        doc.add_paragraph(body)
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def export_pptx(project: Project, structures: Iterable[DocumentStructure]) -> bytes:
    prs = Presentation()
    for structure in sorted(structures, key=lambda s: s.order_index):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = structure.title
        content = structure.content.generated_content if structure.content else "Content pending."
        placeholder = slide.shapes.placeholders[1]
        placeholder.text = content
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


